import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import base64
from datetime import datetime

st.set_page_config(page_title="Sumadhura Ticket Report Generator", layout="wide")

st.markdown("## Sumadhura Ticket Report Generator")
st.caption("Upload the monthly XLSX. Get a CEO-ready PPTX and an interactive HTML deck.")

# Sidebar: branding
with st.sidebar:
    st.header("Branding")
    brand_red = st.color_picker("Primary (Red)", "#e2282a")
    brand_orange = st.color_picker("Accent (Orange)", "#eb7e27")
    brand_green = st.color_picker("Accent (Green)", "#1a9d4a")
    logo_file = st.file_uploader("Logo (PNG)", type=["png"])
    st.divider()
    st.header("Data Settings")
    header_row = st.number_input("Header row (Excel numbering)", min_value=1, value=14, step=1)
    date_col = st.text_input("Date column name", value="Date Reported")
    group_threshold = st.slider("Group small pie slices into 'Other' (min % to label)", 0.0, 10.0, 3.0, 0.5)

uploaded = st.file_uploader("Upload Ticket_Details_Report.xlsx", type=["xlsx"], accept_multiple_files=False)

def hex_to_rgbcolor(hex_str):
    hex_str = hex_str.strip().lstrip('#')
    r = int(hex_str[0:2], 16); g = int(hex_str[2:4], 16); b = int(hex_str[4:6], 16)
    return (r, g, b)

def compute_metrics(df, date_col):
    total_tickets = len(df)
    total_conversations = df['Total Conversations'].sum() if 'Total Conversations' in df.columns else 0
    total_replies = df['Total Replies'].sum() if 'Total Replies' in df.columns else 0
    if date_col not in df.columns:
        raise KeyError(f"Date column '{date_col}' not found. Available columns: {list(df.columns)}")
    df['_date_only'] = pd.to_datetime(df[date_col], errors='coerce').dt.date
    cat_col = 'Category L1(Response)'
    qtype_col = 'Query Type(Response)'
    category_counts = df[cat_col].value_counts().dropna() if cat_col in df.columns else pd.Series(dtype=int)
    query_type_counts = df[qtype_col].value_counts().dropna() if qtype_col in df.columns else pd.Series(dtype=int)
    tickets_per_day = df.groupby('_date_only').size()
    return dict(
        total_tickets=int(total_tickets),
        total_conversations=int(total_conversations),
        total_replies=int(total_replies),
        category_counts=category_counts,
        query_type_counts=query_type_counts,
        tickets_per_day=tickets_per_day
    )

def donut_fig(series, brand_hex, group_threshold=0.03):
    total = series.sum() if len(series)>0 else 0
    fig = plt.figure(figsize=(8,8))
    if total == 0:
        plt.text(0.5,0.5,"No data", ha='center', va='center', fontsize=16); plt.axis('off')
        return fig, []
    large = series[series/total >= group_threshold]
    small = series[series/total < group_threshold]
    grouped = large.copy()
    if not small.empty:
        grouped['Other'] = small.sum()
    wedges, texts, autotexts = plt.pie(
        grouped.values,
        labels=None,
        autopct=lambda pct: ('%1.1f%%' % pct) if pct >= group_threshold*100 else '',
        startangle=140,
        pctdistance=0.75,
        textprops={'fontsize': 12},
        wedgeprops={'linewidth': 1, 'edgecolor': 'white'},
        colors=[brand_hex['red'], brand_hex['orange'], brand_hex['green'], '#888888', '#cccccc']
    )
    centre_circle = plt.Circle((0,0), 0.45, fc='white'); fig.gca().add_artist(centre_circle)
    plt.title('Ticket Category Breakdown (L1)', fontsize=16)
    legend_labels = [f'{name} ({val})' for name, val in zip(grouped.index, grouped.values)]
    plt.legend(wedges, legend_labels, title='Category (count)', loc='center left', bbox_to_anchor=(1, 0.5), fontsize=11, title_fontsize=12)
    plt.tight_layout()
    return fig, legend_labels

def bar_fig(series, title):
    fig = plt.figure(figsize=(10,6))
    if len(series)==0:
        plt.text(0.5,0.5,"No data", ha='center', va='center', fontsize=16); plt.axis('off')
    else:
        series.plot(kind='bar')
        plt.title(title, fontsize=14); plt.xlabel(''); plt.ylabel('Count')
        plt.xticks(rotation=45, ha='right')
    plt.tight_layout(); return fig

def line_fig(series, title):
    fig = plt.figure(figsize=(10,6))
    if len(series)==0:
        plt.text(0.5,0.5,"No data", ha='center', va='center', fontsize=16); plt.axis('off')
    else:
        series.plot(kind='line', marker='o')
        plt.title(title, fontsize=14); plt.xlabel('Date'); plt.ylabel('Count')
        plt.xticks(rotation=45, ha='right'); plt.grid(True)
    plt.tight_layout(); return fig

def fig_to_png_bytes(fig):
    buf = BytesIO(); fig.savefig(buf, format='png', dpi=180, bbox_inches='tight'); buf.seek(0); return buf

def build_pptx(metrics, logo_png_bytes, brand_hex):
    prs = Presentation()
    red = RGBColor(*hex_to_rgbcolor(brand_hex['red']))
    orange = RGBColor(*hex_to_rgbcolor(brand_hex['orange']))
    green = RGBColor(*hex_to_rgbcolor(brand_hex['green']))

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Simplify360 Ticket Report"
    subtitle = slide.placeholders[1]; subtitle.text = "Client: Sumadhura Group"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(44); subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    if logo_png_bytes:
        slide.shapes.add_picture(logo_png_bytes, Inches(8), Inches(0.3), width=Inches(1.5))

    # KPI slide
    s2 = prs.slides.add_slide(prs.slide_layouts[5])
    s2.shapes.title.text = "Key Metrics"
    s2.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(8), Inches(1)
    vals = [("Total Tickets", metrics['total_tickets'], green),
            ("Total Conversations", metrics['total_conversations'], orange),
            ("Total Replies", metrics['total_replies'], red)]
    for i,(label,val,color) in enumerate(vals):
        tb = s2.shapes.add_textbox(left, top + Inches(i*1.2), width, height)
        p = tb.text_frame.paragraphs[0]; p.text = f"{label}: {val}"
        p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = color
    return prs

def build_inline_html(metrics, logo_png_bytes, brand_hex, donut_png, bar_png, line_png):
    def b64(b): return base64.b64encode(b).decode('utf-8')
    logo_b64 = b64(logo_png_bytes) if logo_png_bytes else ''
    html = f"""<!doctype html>
<html>
<head>
  <meta charset="utf-8">
  <title>Sumadhura | Simplify360 Ticket Report</title>
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <link rel="stylesheet" href="https://unpkg.com/reveal.js@5/dist/reveal.css">
  <link rel="stylesheet" href="https://unpkg.com/reveal.js@5/dist/theme/black.css" id="theme">
  <style>
    :root {{
      --brand-red: {brand_hex['red']}; --brand-orange: {brand_hex['orange']}; --brand-green: {brand_hex['green']};
    }}
    .reveal {{ font-size: 28px; }}
    .brand-bar {{ height: 8px; display: grid; grid-template-columns: 1fr 1fr 1fr; }}
    .brand-bar > div:nth-child(1) {{ background: var(--brand-red); }}
    .brand-bar > div:nth-child(2) {{ background: var(--brand-orange); }}
    .brand-bar > div:nth-child(3) {{ background: var(--brand-green); }}
    .kpi {{ display: grid; grid-template-columns: repeat(3, 1fr); gap: 24px; margin-top: 24px; }}
    .kpi div {{ background: rgba(255,255,255,0.06); border-radius: 16px; padding: 20px; text-align: center; box-shadow: 0 6px 18px rgba(0,0,0,0.2); }}
    .kpi span {{ display:block; color:#bbb; font-size:22px; }}
    .kpi strong {{ display:block; font-size:56px; margin-top:8px; font-weight:800; color: var(--brand-green); }}
    .logo {{ position:absolute; top:16px; right:16px; width:140px; }}
    .caption {{ color:#bbb; font-size:20px; margin-top:10px; }}
    img.chart {{ border-radius: 14px; box-shadow: 0 6px 18px rgba(0,0,0,0.25); }}
    .reveal .slides section {{ overflow-y: auto !important; }}
    .chart-wrap {{ max-height: 75vh; overflow: auto; display: grid; place-items: center; }}
    img.chart {{ max-width: 95vw; max-height: 70vh; object-fit: contain; }}
  </style>
</head>
<body>
<div class="reveal"><div class="slides">
  <section data-background-color="#000">
    {('<img class="logo" src="data:image/png;base64,'+logo_b64+'" />') if logo_b64 else ''}
    <h1> Simplify360 Ticket Report </h1>
    <h3> Client: Sumadhura Group </h3>
    <p style="opacity:.7">Updated: {datetime.today().strftime("%b %d, %Y")}</p>
    <div class="brand-bar"><div></div><div></div><div></div></div>
    <p class="caption">Auto-generated deck — optimized for CEO review</p>
  </section>
  <section>
    <h2>Key Metrics</h2>
    <div class="kpi">
      <div><span>Total Tickets</span><strong>{metrics['total_tickets']}</strong></div>
      <div><span>Total Conversations</span><strong>{metrics['total_conversations']}</strong></div>
      <div><span>Total Replies</span><strong>{metrics['total_replies']}</strong></div>
    </div>
    <div class="brand-bar" style="margin-top:32px"><div></div><div></div><div></div></div>
  </section>
  <section>
    <h2>Ticket Category Breakdown (L1)</h2>
    <div class="chart-wrap"><img class="chart" src="data:image/png;base64,{base64.b64encode(donut_png).decode('utf-8')}" alt="Donut"></div>
    <p class="caption">Small categories grouped into “Other” for readability</p>
  </section>
  <section>
    <h2>Query Type Breakdown</h2>
    <div class="chart-wrap"><img class="chart" src="data:image/png;base64,{base64.b64encode(bar_png).decode('utf-8')}" alt="Bar"></div>
  </section>
  <section>
    <h2>Tickets Over Time</h2>
    <div class="chart-wrap"><img class="chart" src="data:image/png;base64,{base64.b64encode(line_png).decode('utf-8')}" alt="Line"></div>
  </section>
</div></div>
<script src="https://unpkg.com/reveal.js@5/dist/reveal.js"></script>
<script>Reveal.initialize({hash:true, slideNumber:true, transition:'fade'});</script>
</body></html>"""
    return html.encode('utf-8')

if uploaded:
    try:
        # Load sheet with chosen header row (Excel is 1-based)
        df = pd.read_excel(uploaded, sheet_name='Data', header=header_row-1)
        metrics = compute_metrics(df, date_col)

        # Brand colors
        brand_hex = {'red': brand_red, 'orange': brand_orange, 'green': brand_green}

        # Charts
        donut, _ = donut_fig(metrics['category_counts'], brand_hex, group_threshold=group_threshold/100.0)
        bar = bar_fig(metrics['query_type_counts'], "Query Type Breakdown")
        line = line_fig(metrics['tickets_per_day'], "Tickets Over Time")

        c1, c2, c3 = st.columns(3)
        with c1: st.pyplot(donut, use_container_width=True)
        with c2: st.pyplot(bar, use_container_width=True)
        with c3: st.pyplot(line, use_container_width=True)

        # Collect logo bytes if provided
        logo_bytes = logo_file.read() if logo_file is not None else None

        # Create PPTX
        prs = build_pptx(metrics, BytesIO(logo_bytes) if logo_bytes else None, brand_hex)
        # Attach chart images to PPTX slides
        def add_chart_slide(prs, title, fig):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = title
            png = fig_to_png_bytes(fig)
            png.seek(0)
            slide.shapes.add_picture(png, Inches(0.5), Inches(1.5), width=Inches(8))

        add_chart_slide(prs, "Ticket Category Breakdown (L1)", donut)
        add_chart_slide(prs, "Query Type Breakdown", bar)
        add_chart_slide(prs, "Tickets Over Time", line)

        pptx_buf = BytesIO(); prs.save(pptx_buf); pptx_buf.seek(0)

        # Build inline HTML
        donut_png = fig_to_png_bytes(donut).getvalue()
        bar_png = fig_to_png_bytes(bar).getvalue()
        line_png = fig_to_png_bytes(line).getvalue()
        html_bytes = build_inline_html(metrics, logo_bytes, brand_hex, donut_png, bar_png, line_png)

        st.success("Report generated! Use the download buttons below.")
        st.download_button("⬇️ Download PPTX", data=pptx_buf.getvalue(), file_name="Sumadhura_Ticket_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        st.download_button("⬇️ Download Interactive HTML", data=html_bytes, file_name="Sumadhura_Ticket_Report_inline.html", mime="text/html")

    except Exception as e:
        st.error(f"Error: {e}")
else:
    st.info("Upload the monthly XLSX on this page to generate the report.")
